from io import BytesIO
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pypdf import PdfReader


def load_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    pages = []

    for page in reader.pages:
        pages.append(page.extract_text() or "")

    return "\n\n".join(page for page in pages if page)


def load_docx(file_bytes: bytes) -> str:
    document = Document(BytesIO(file_bytes))
    paragraphs = [paragraph.text for paragraph in document.paragraphs]

    return "\n".join(paragraph for paragraph in paragraphs if paragraph.strip())


def load_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))
    slide_text_blocks = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = ""
        text_lines = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines = [
                    line.strip()
                    for line in shape.text.splitlines()
                    if line.strip()
                ]
                if not lines:
                    continue

                if (
                    shape.is_placeholder
                    and shape.placeholder_format.type
                    in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
                    and not title
                ):
                    title = lines[0]

                text_lines.extend(lines)

        if not title and text_lines:
            title = text_lines[0]

        content_lines = []
        seen_lines = set()

        for line in text_lines:
            if title and line.strip().lower() == title.strip().lower():
                continue

            normalized_line = line.lower()
            if normalized_line in seen_lines:
                continue

            content_lines.append(line)
            seen_lines.add(normalized_line)

        if title or content_lines:
            block_lines = [f"[Slide {slide_index}]"]

            if title:
                block_lines.append(f"Title: {title}")

            block_lines.append("Content:")
            block_lines.extend(f"- {line}" for line in content_lines)

            slide_text_blocks.append("\n".join(block_lines))

    return "\n\n".join(slide_text_blocks)


def load_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("utf-8", errors="replace")


def load_uploaded_documents(uploaded_files) -> list[dict]:
    documents = []

    for uploaded_file in uploaded_files:
        file_name = uploaded_file.name
        suffix = Path(file_name).suffix.lower()
        file_bytes = uploaded_file.getvalue()

        if suffix == ".pdf":
            extracted_text = load_pdf(file_bytes)
        elif suffix == ".docx":
            extracted_text = load_docx(file_bytes)
        elif suffix == ".pptx":
            extracted_text = load_pptx(file_bytes)
        elif suffix == ".txt":
            extracted_text = load_txt(file_bytes)
        else:
            continue

        documents.append(
            {
                "source": file_name,
                "text": extracted_text,
            }
        )

    return documents
